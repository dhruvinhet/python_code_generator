import os
from PyPDF2 import PdfReader
import numpy as np
import google.generativeai as genai
import sqlite3
import json
import faiss
import logging

from sentence_transformers import SentenceTransformer
from docx import Document
from pptx import Presentation
from flask import current_app as app

# Import the local embedding model from config
from .config import (
    local_embedding_model_instance, 
    GEMINI_EMBEDDING_MODEL_API, 
    LOCAL_EMBEDDING_MODEL_NAME
)

DB_PATH = 'quiz_data.db'
FAISS_INDEX_DIR = 'faiss_indexes'
os.makedirs(FAISS_INDEX_DIR, exist_ok=True)

def extract_text_from_document(file_path):
    """Extracts text from PDF, DOCX, or PPTX files, returning a list of page/slide texts."""
    ext = os.path.splitext(file_path)[1].lower()
    text_pages = []
    try:
        if ext == '.pdf':
            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_pages.append(page_text.strip())
        elif ext == '.docx':
            doc = Document(file_path)
            PAGE_SIZE = 10
            current_page = []
            for para in doc.paragraphs:
                if para.text.strip():
                    current_page.append(para.text.strip())
                    if len(current_page) >= PAGE_SIZE:
                        text_pages.append('\n'.join(current_page))
                        current_page = []
            if current_page:
                text_pages.append('\n'.join(current_page))
        elif ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                slide_text = '\n'.join([shape.text.strip() for shape in slide.shapes if shape.has_text_frame and shape.text.strip()])
                if slide_text:
                    text_pages.append(slide_text)
        else:
            print(f"Unsupported file type: {ext}")
            return None
        
        if not text_pages:
            return None
        
        print(f"\n--- Extracted {len(text_pages)} pages/slides from {ext} file ---")
        print("--- First page/slide (first 200 chars):", text_pages[0][:200] if text_pages else "No content ---")
        return clean_text(text_pages)
    except Exception as e:
        print(f"Error extracting text from document: {e}")
        return None

def clean_text(pages):
    cleaned_pages = []
    irrelevant_keywords = [
        "course structure", "grading", "marks", "university", "syllabus", 
        "table of contents", "index", "glossary", "references", "bibliography"
    ]

    for page_text in pages:
        import re
        # Check if the page contains a high percentage of irrelevant keywords
        lower_page_text = page_text.lower()
        if sum(keyword in lower_page_text for keyword in irrelevant_keywords) > 2:
            continue

        # Remove common irrelevant patterns
        patterns_to_remove = [
            r"University of [\w\s]+",
            r"Marks: \d+",
            r"Page \d+ of \d+",
            r"^\s*table of contents\s*$",
            r"^\s*index\s*$",
            r"^\s*glossary\s*$",
            r"^\s*references\s*$",
            r"^\s*bibliography\s*$",
            r"^(unit|chapter|section)\s+\d+",
            r"^\s*[A-Z\s]{5,}\s*$",
        ]
        
        cleaned_page_text = page_text
        for pattern in patterns_to_remove:
            cleaned_page_text = re.sub(pattern, "", cleaned_page_text, flags=re.IGNORECASE | re.MULTILINE)

        # Remove identified headers and footers
        lines = cleaned_page_text.split('\n')
        cleaned_lines = [line for line in lines if len(line.strip()) > 20] # Keep only longer lines
        
        cleaned_page_text = '\n'.join(cleaned_lines)
        
        # Further clean up whitespace
        cleaned_page_text = re.sub(r'\n{3,}', '\n\n', cleaned_page_text)
        
        if cleaned_page_text.strip():
            cleaned_pages.append(cleaned_page_text.strip())
            
    return cleaned_pages


def chunk_text(pages, chunk_size=250, chunk_overlap=25):
    """
    Splits joined text from pages into smaller, overlapping chunks.
    """
    text = '\n\n'.join(pages)  # Join pages for chunking
    chunks = []
    if not text:
        return chunks
    
    words = text.split()
    current_chunk = []
    current_len = 0

    for word in words:
        if current_len + len(word) + 1 > chunk_size and current_chunk:
            chunks.append(" ".join(current_chunk))
            overlap_words = int(chunk_overlap / (len(current_chunk[-1]) + 1) * len(current_chunk)) if current_chunk else 0
            current_chunk = current_chunk[-overlap_words:]
            current_len = sum(len(w) + 1 for w in current_chunk) - 1 if current_chunk else 0

        current_chunk.append(word)
        current_len += len(word) + 1

    if current_chunk:
        chunks.append(" ".join(current_chunk))
    
    print(f"\n--- Number of Chunks Created: {len(chunks)} ---")
    print("--- First Chunk (first 200 chars):", chunks[0][:200] if chunks else "No chunks ---")
    return chunks

def get_embeddings(texts):
    """Generates embeddings for a list of texts using a local model or Gemini's embedding model."""
    try:
        if local_embedding_model_instance:
            print(f"Using local embedding model: {LOCAL_EMBEDDING_MODEL_NAME}")
            embeddings = local_embedding_model_instance.encode(texts, convert_to_numpy=True)
        else:
            print(f"Using Gemini embedding model: {GEMINI_EMBEDDING_MODEL_API}")
            response = genai.embed_content(
                model=GEMINI_EMBEDDING_MODEL_API,
                content=texts,
                task_type="RETRIEVAL_DOCUMENT"
            )
            embeddings = np.array(response['embedding']).astype('float32')
        
        print(f"Generated {len(texts)} embeddings (dimension: {embeddings.shape[1]})")
        return embeddings
    except Exception as e:
        print(f"Error generating embeddings: {e}")
        if hasattr(e, 'status_code') and e.status_code == 403:
            print("Google API Key/Service Blocked Error: Check your Google Cloud Console for API Key restrictions and enabled APIs.")
        return None


def init_db():
    """Initializes the database schema."""
    logging.info(f"Initializing database at: {DB_PATH}") # Added logging
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    
    # Ensure all tables exist
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS documents (
            id TEXT PRIMARY KEY,
            original_filename TEXT NOT NULL,
            faiss_index_path TEXT NOT NULL,
            chunks_path TEXT NOT NULL,
            pages_path TEXT NOT NULL
        )
    """)
    
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS quizzes (
            quiz_id TEXT PRIMARY KEY,
            document_id TEXT,
            quiz_type TEXT,
            num_questions INTEGER,
            quiz_data TEXT,
            quiz_results TEXT DEFAULT '[]',
            timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY(document_id) REFERENCES documents(id)
        )
    """)
    
    # Add migration logic for existing quizzes table
    try:
        # Check if quiz_results column exists
        cursor.execute("PRAGMA table_info(quizzes)")
        columns = [column[1] for column in cursor.fetchall()]
        
        if 'quiz_results' not in columns:
            cursor.execute("ALTER TABLE quizzes ADD COLUMN quiz_results TEXT DEFAULT '[]'")
            logging.info("Added quiz_results column to quizzes table")
            
        if 'timestamp' not in columns:
            cursor.execute("ALTER TABLE quizzes ADD COLUMN timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP")
            logging.info("Added timestamp column to quizzes table")
            
    except sqlite3.Error as e:
        logging.warning(f"Migration warning: {e}")
    
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS quiz_submissions (
            quiz_results_id INTEGER PRIMARY KEY AUTOINCREMENT,
            quiz_id TEXT,
            user_answers TEXT,
            quiz_results TEXT,
            FOREIGN KEY(quiz_id) REFERENCES quizzes(quiz_id)
        )
    """)
    
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS chat_sessions (
            chat_id TEXT PRIMARY KEY,
            document_id TEXT,
            quiz_type TEXT,
            messages TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY(document_id) REFERENCES documents(id)
        )
    """)

    conn.commit()
    conn.close()

def save_document_data(doc_id, original_filename, faiss_index, chunks, pages):
    init_db()
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()

    index_path = os.path.join(FAISS_INDEX_DIR, f"{doc_id}.faiss")
    logging.info(f"Attempting to save FAISS index to: {index_path}")
    try:
        faiss.write_index(faiss_index, index_path)
        logging.info(f"Successfully saved FAISS index to: {index_path}")
    except Exception as e:
        logging.error(f"Error saving FAISS index to {index_path}: {e}")
        # Optionally re-raise or handle more gracefully

    pages_path = os.path.join(FAISS_INDEX_DIR, f"{doc_id}_pages.json")
    with open(pages_path, 'w', encoding='utf-8') as f:
        json.dump(pages, f)
    logging.info(f"Successfully saved pages to: {pages_path}")
    
    chunks_path = os.path.join(FAISS_INDEX_DIR, f"{doc_id}_chunks.json")
    with open(chunks_path, 'w', encoding='utf-8') as f:
        json.dump(chunks, f)
    logging.info(f"Successfully saved chunks to: {chunks_path}")

    cursor.execute('''
        INSERT OR REPLACE INTO documents (id, original_filename, faiss_index_path, pages_path, chunks_path)
        VALUES (?, ?, ?, ?, ?)
    ''', (doc_id, original_filename, index_path, pages_path, chunks_path))
    conn.commit()
    conn.close()

def load_document_data(doc_id):
    init_db()
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute('SELECT original_filename, faiss_index_path, pages_path, chunks_path FROM documents WHERE id = ?', (doc_id,))
    row = cursor.fetchone()
    conn.close()

    if row:
        original_filename, index_path, pages_path, chunks_path = row
        faiss_index = faiss.read_index(index_path)
        
        with open(pages_path, 'r', encoding='utf-8') as f:
            pages = json.load(f)
        with open(chunks_path, 'r', encoding='utf-8') as f:
            chunks = json.load(f)

        return {
            "original_filename": original_filename,
            "pages": pages,
            "chunks": chunks,
            "faiss_index": faiss_index
        }
    return None 

def get_all_documents_meta():
    init_db()
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute('SELECT id, original_filename FROM documents ORDER BY original_filename ASC')
    docs = [{"id": row[0], "filename": row[1]} for row in cursor.fetchall()]
    conn.close()
    return docs


def save_quiz_to_db(quiz_id, document_id, quiz_type, num_questions, quiz_data):
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    try:
        # Provide an empty JSON string for quiz_results to satisfy the NOT NULL constraint
        cursor.execute('''
            INSERT INTO quizzes (quiz_id, document_id, quiz_type, num_questions, quiz_data, quiz_results)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (quiz_id, document_id, quiz_type, num_questions, json.dumps(quiz_data), json.dumps([])))
        conn.commit()
    except sqlite3.Error as e:
        print(f"Database error in save_quiz_to_db: {e}")
        conn.rollback()
        raise
    finally:
        conn.close()

def get_quizzes_for_document(document_id):
    init_db()
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute('SELECT quiz_id, quiz_type, num_questions, quiz_data, timestamp FROM quizzes WHERE document_id = ? ORDER BY timestamp DESC', (document_id,))
    quizzes = []
    for row in cursor.fetchall():
        quiz_id, quiz_type, num_questions, quiz_data_json, timestamp = row
        quizzes.append({
            "quiz_id": quiz_id,
            "quiz_type": quiz_type,
            "num_questions": num_questions,
            "quiz_data": json.loads(quiz_data_json),
            "timestamp": timestamp
        })
    conn.close()
    return quizzes


def get_questions_for_document(document_id):
    quizzes = get_quizzes_for_document(document_id)
    questions = []
    for quiz in quizzes:
        for question in quiz['quiz_data']:
            questions.append(question['question'])
    return questions


# in backend/document_processor.py

def save_chat_history_to_db(chat_id, document_id, quiz_type, messages):
    """Saves a chat session's history to the database."""
    conn = None
    try:
        logging.info(f"Attempting to save chat history: chat_id={chat_id}, document_id={document_id}, quiz_type={quiz_type}")
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        
        # Check if the chat session already exists
        cursor.execute("SELECT chat_id FROM chat_sessions WHERE chat_id = ?", (chat_id,))
        exists = cursor.fetchone()

        if exists:
            # If the session exists, update the messages
            cursor.execute("""
                UPDATE chat_sessions
                SET messages = ?, quiz_type = ?
                WHERE chat_id = ?
            """, (json.dumps(messages), quiz_type, chat_id))
            logging.info(f"Updated chat session {chat_id} in the database.")
        else:
            # If it's a new session, insert a new record
            cursor.execute("""
                INSERT INTO chat_sessions (chat_id, document_id, quiz_type, messages)
                VALUES (?, ?, ?, ?)
            """, (chat_id, document_id, quiz_type, json.dumps(messages)))
            logging.info(f"Saved new chat session {chat_id} to the database.")
            
        conn.commit()
        logging.info(f"Changes committed for chat_id: {chat_id}")
    except sqlite3.Error as e:
        logging.error(f"SQLite error during chat history save for chat_id {chat_id}: {e}")
        if conn:
            conn.rollback()  # Rollback changes on error
    except Exception as e:
        logging.error(f"Unexpected error during chat history save for chat_id {chat_id}: {e}")
        if conn:
            conn.rollback() # Rollback changes on error
    finally:
        if conn:
            conn.close()

def get_chat_history(chat_id):
    init_db()
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute('SELECT messages FROM chat_sessions WHERE chat_id = ?', (chat_id,))
    row = cursor.fetchone()
    conn.close()
    if row:
        return json.loads(row[0])
    return []

def get_all_chat_sessions_meta():
    """Retrieves metadata for all chat sessions."""
    conn = None
    sessions = []
    try:
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        
        logging.info(f"Attempting to fetch chat sessions from database at: {DB_PATH}")

        # Use a LEFT JOIN to ensure all chat sessions are returned,
        # even if a matching document record is missing.
        cursor.execute("""
            SELECT 
                cs.chat_id, 
                cs.quiz_type, 
                cs.document_id, 
                d.original_filename
            FROM chat_sessions cs
            LEFT JOIN documents d ON cs.document_id = d.id
            ORDER BY cs.created_at DESC
        """)
        
        rows = cursor.fetchall()
        logging.info(f"Raw rows fetched from chat_sessions: {rows}") # Added logging
        
        if not rows:
            logging.warning("No chat sessions found in the database.")
            return []

        for row in rows:
            filename = row[3] if row[3] else "Document Not Found"
            
            sessions.append({
                "chat_id": row[0],
                "quiz_type": row[1],
                "document_id": row[2],
                "document_filename": filename
            })
            
        logging.info(f"Successfully retrieved {len(sessions)} chat sessions.")
        
    except sqlite3.Error as e:
        logging.error(f"SQLite error while fetching chat sessions: {e}")
        sessions = [] # Return empty list on error
    finally:
        if conn:
            conn.close()
            
    return sessions

# In backend/document_processor.py

def get_quiz_results_from_db(quiz_id):
    init_db()
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    # The SQL query remains the same
    cursor.execute('SELECT user_answers, quiz_results, timestamp FROM quiz_submissions WHERE quiz_id = ? ORDER BY timestamp DESC LIMIT 1', (quiz_id,))
    row = cursor.fetchone()
    conn.close()
    if row:
        user_answers = json.loads(row[0])
        quiz_results_json_string = row[1]
        timestamp = row[2]
        
        # Parse the JSON string into a Python object
        try:
            quiz_results = json.loads(quiz_results_json_string)
        except json.JSONDecodeError as e:
            print(f"Failed to decode quiz_results from DB: {e}")
            quiz_results = [] # Provide a safe fallback

        return {"user_answers": user_answers, "quiz_results": quiz_results, "timestamp": timestamp}
    return None

def save_quiz_results_to_db(quiz_id, user_answers, quiz_results):
    init_db()
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    try:
        app.logger.info(f"Saving quiz results for quiz_id: {quiz_id}")
        cursor.execute('''
            INSERT INTO quiz_submissions (quiz_id, user_answers, quiz_results)
            VALUES (?, ?, ?)
        ''', (quiz_id, json.dumps(user_answers), json.dumps(quiz_results)))
        conn.commit()
        app.logger.info(f"Successfully saved quiz results for quiz_id: {quiz_id}")
    except sqlite3.Error as e:
        app.logger.error(f"Database error in save_quiz_results_to_db: {e}")
        conn.rollback()
        raise
    finally:
        conn.close()
